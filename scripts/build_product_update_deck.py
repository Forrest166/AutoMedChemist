from __future__ import annotations

import json
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
PPTX_OUT = ROOT / "AutoMedChemist_Product_Update.pptx"
PDF_OUT = ROOT / "AutoMedChemist_Product_Update.pdf"
PREVIEW_DIR = ROOT / "docs" / "product_update_previews"

INK = RGBColor(23, 32, 42)
MUTED = RGBColor(82, 97, 111)
TEAL = RGBColor(15, 118, 110)
GREEN = RGBColor(31, 122, 82)
AMBER = RGBColor(180, 83, 9)
CORAL = RGBColor(190, 65, 55)
PAPER = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)


def read_json(path: Path) -> dict:
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8")) or {}
    except Exception:
        return {}
    return data if isinstance(data, dict) else {}


def compact_status(value: object) -> str:
    text = str(value or "-")
    return {
        "attention_required": "attention",
        "review_required": "review",
        "needs_attention": "attention",
        "awaiting_exact_results": "awaiting",
    }.get(text, text.replace("_", " "))


def text_box(slide, x, y, w, h, text, size=20, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullet_list(slide, x, y, w, items, size=17, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(4.5))
    frame = box.text_frame
    frame.clear()
    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.level = 0
        paragraph.space_after = Pt(8)
    return box


def add_metric(slide, x, y, label, value, color=TEAL):
    text_box(slide, x, y, 2.0, 0.35, label.upper(), size=8, bold=True, color=MUTED)
    text_box(slide, x, y + 0.28, 2.0, 0.5, str(value), size=24, bold=True, color=color)


def add_bar(slide, x, y, w, label, value, max_value, color):
    text_box(slide, x, y - 0.02, 2.1, 0.25, label, size=10, bold=True, color=INK)
    bg = slide.shapes.add_shape(1, Inches(x + 2.3), Inches(y), Inches(w), Inches(0.16))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(222, 228, 234)
    bg.line.fill.background()
    width = 0 if max_value <= 0 else max(0.05, min(w, w * float(value) / float(max_value)))
    fg = slide.shapes.add_shape(1, Inches(x + 2.3), Inches(y), Inches(width), Inches(0.16))
    fg.fill.solid()
    fg.fill.fore_color.rgb = color
    fg.line.fill.background()
    text_box(slide, x + 2.3 + w + 0.12, y - 0.08, 0.8, 0.3, str(value), size=11, bold=True, color=INK)


def add_title(slide, title, kicker="AutoMedChemist product update"):
    text_box(slide, 0.65, 0.42, 4.5, 0.3, kicker.upper(), size=8, bold=True, color=TEAL)
    text_box(slide, 0.65, 0.72, 8.6, 0.75, title, size=25, bold=True, color=INK)


def build_slide_data() -> list[dict]:
    readiness = read_json(ROOT / "data/projects/demo/promotion_readiness_packet.json")
    dashboard = read_json(ROOT / "data/releases/production_dashboard_snapshot.json")
    production_ci = read_json(ROOT / "data/releases/production_ci_report.json")
    production_smoke = read_json(ROOT / "data/releases/release_smoke_checklist_production.json")
    native_smoke = read_json(ROOT / "data/releases/native_shell_smoke.json")
    quality = read_json(ROOT / "data/releases/native_ui_quality_report.json")
    regression = read_json(ROOT / "data/releases/native_ui_regression_snapshot.json")
    portable = read_json(ROOT / "data/releases/native_portable_package_manifest.json")
    data_foundation = read_json(ROOT / "data/substituents/data_foundation_report.json")
    board = read_json(ROOT / "data/projects/demo/candidate_review_board.json")
    analytics = read_json(ROOT / "data/projects/demo/candidate_review_analytics.json")
    decision = read_json(ROOT / "data/projects/demo/candidate_decision_packet.json")
    evidence_quality = read_json(ROOT / "data/projects/demo/evidence_quality_scorecard.json")
    reviewer_ops = read_json(ROOT / "data/projects/demo/reviewer_operations.json")
    lineage = read_json(ROOT / "data/projects/demo/baseline_lineage_compare.json")
    trend_summary = read_json(ROOT / "data/releases/operator_trend_summary.json")
    trend_charts = read_json(ROOT / "data/releases/operator_trend_charts.json")
    component_locator = read_json(ROOT / "data/projects/demo/candidate_component_structure_locator.json")
    site_calibration = read_json(ROOT / "data/projects/demo/site_detection_calibration_queue.json")
    rgroup_replay = read_json(ROOT / "data/substituents/rgroup_admission_sandbox_impact_replay.json")
    reviewer_cockpit = read_json(ROOT / "data/projects/demo/reviewer_cockpit.json")
    totals = data_foundation.get("totals") or {}
    flag_counts = evidence_quality.get("flag_counts") or {}
    cockpit_lanes = reviewer_cockpit.get("lane_counts") or {}
    cockpit_lane_summary = ", ".join(
        f"{key}={value}" for key, value in cockpit_lanes.items()
    ) or "no lanes yet"
    return [
        {
            "title": "AutoMedChemist is now a native review workbench",
            "subtitle": "The desktop shell connects local candidate generation, review analytics, evidence quality, baseline lineage, and production gates without opening a browser.",
            "bullets": [
                f"Native shell smoke: {native_smoke.get('status') or 'not run'}",
                f"High-DPI UI quality report: {quality.get('status') or 'not run'}",
                f"Native regression snapshot: {regression.get('status') or 'not run'}",
                "User-facing path: AutoMedChemist.exe in the project root.",
            ],
        },
        {
            "title": "Release gates are green after the advanced review pass",
            "metrics": [
                ("Dashboard", dashboard.get("row_count") or 0),
                ("CI steps", production_ci.get("step_count") or 0),
                ("Smoke", len(production_smoke.get("checks") or [])),
                ("Assets", totals.get("asset_count") or 0),
            ],
            "bullets": [
                f"Production dashboard: {compact_status(dashboard.get('status'))}, fail={dashboard.get('fail_count') or 0}, warn={dashboard.get('warn_count') or 0}.",
                f"Production CI: {compact_status(production_ci.get('status'))} with {len(production_ci.get('failed_steps') or [])} failed steps.",
                f"Data foundation: {totals.get('missing_asset_count') or 0} missing assets and {totals.get('warning_count') or 0} warnings.",
                "All new artifacts remain local decision support; external operational workflows stay blocked.",
            ],
        },
        {
            "title": "Advanced local review layer is now wired",
            "metrics": [
                ("Locators", component_locator.get("linked_component_count") or 0),
                ("Calibration", site_calibration.get("queue_count") or 0),
                ("R-group replay", rgroup_replay.get("row_count") or 0),
                ("Cockpit", reviewer_cockpit.get("row_count") or 0),
            ],
            "bullets": [
                "Candidate score components now route back to the same 2D before/after structure view and highlight detail.",
                f"Site detection calibration queue tracks {site_calibration.get('low_confidence_count') or 0} low-confidence rows as local parser fixtures.",
                f"R-group admission replay keeps {rgroup_replay.get('rollback_ready_count') or 0} rollback-ready source rows with production scoring writes blocked.",
                f"Reviewer Cockpit merges reason audit, closure, and remediation lanes: {cockpit_lane_summary}.",
            ],
        },
        {
            "title": "Review analytics now drives the board",
            "metrics": [
                ("Board rows", board.get("row_count") or 0),
                ("Focused", board.get("focused_row_count") or 0),
                ("Pending", board.get("pending_local_review_count") or 0),
                ("Backlog", analytics.get("pending_backlog_count") or 0),
            ],
            "bullets": [
                "Native analytics rows carry filter hints for site class, non-clear risk, reviewer workload, attention, and all-row views.",
                "Selecting a review analytics card filters the Candidate Review table and refreshes evidence context.",
                "Batch review status updates still write only local review decisions, not execution actions.",
            ],
        },
        {
            "title": "Evidence quality became an explicit scorecard",
            "metrics": [
                ("Rows", evidence_quality.get("row_count") or 0),
                ("Attention", evidence_quality.get("attention_count") or 0),
                ("Watch", evidence_quality.get("watch_count") or 0),
                ("Thin", flag_counts.get("thin_mmp_sar_evidence") or 0),
            ],
            "bullets": [
                "The scorecard separates thin MMP/SAR evidence, contradiction-heavy evidence, stale review rows, and missing baseline context.",
                f"Contradiction-heavy rows: {flag_counts.get('contradiction_heavy_evidence') or 0}; stale review rows: {flag_counts.get('stale_review_row') or 0}.",
                "Compatibility alias artifacts remain available for older candidate_evidence_quality readers.",
            ],
        },
        {
            "title": "Reviewer operations exposes local workload risk",
            "metrics": [
                ("Rows", reviewer_ops.get("row_count") or 0),
                ("Overdue", reviewer_ops.get("pending_overdue_count") or 0),
                ("Repeated", reviewer_ops.get("repeated_defer_reason_count") or 0),
                ("Pending", reviewer_ops.get("workload_pending_count") or 0),
            ],
            "bullets": [
                "Operations rows cover pending SLA bands, repeated deferral reasons, reviewer workload, ledger events, and site-class closure rates.",
                f"Low site-class closure rows: {reviewer_ops.get('low_site_class_closure_count') or 0}; reviewer lanes: {reviewer_ops.get('reviewer_count') or 0}.",
                "This is a routing aid for local review ownership, not an external task feed.",
            ],
        },
        {
            "title": "Baseline lineage explains candidate movement",
            "metrics": [
                ("Candidates", lineage.get("row_count") or 0),
                ("Entered", lineage.get("entered_candidate_count") or 0),
                ("Exited", lineage.get("exited_candidate_count") or 0),
                ("Changed", lineage.get("changed_candidate_count") or 0),
            ],
            "bullets": [
                f"Current compare: {lineage.get('base_baseline_id') or '-'} to {lineage.get('head_baseline_id') or '-'}.",
                "Rows carry entered, exited, changed, or unchanged lineage status plus local rationale.",
                "Legacy candidate_baseline_lineage artifacts are still emitted for compatibility.",
            ],
        },
        {
            "title": "Operator trend previews are now in-app",
            "metrics": [
                ("Cards", trend_summary.get("card_count") or 0),
                ("Attention", trend_summary.get("needs_attention_count") or 0),
                ("Charts", trend_charts.get("chart_count") or 0),
                ("Regression", len(regression.get("checks") or [])),
            ],
            "bullets": [
                "Trend cards now include review analytics, evidence quality, reviewer operations, baseline movement, DB latency, and packet coverage.",
                "Chart generation writes SVG cards plus high-resolution PNG previews for Native Reports.",
                f"Portable package rebuilt with {portable.get('copied_file_count') or 0} files; package status is {compact_status(portable.get('status'))}.",
            ],
        },
        {
            "title": "Next stage: deeper local design intelligence",
            "metrics": [
                ("Ready", readiness.get("readiness_score") or 0),
                ("Decisions", decision.get("decision_count") or 0),
                ("Defer", decision.get("defer_count") or 0),
                ("Watch", decision.get("watch_count") or 0),
            ],
            "bullets": [
                "Upgrade 2D candidate explanation from linked panels to atom-level visual legends for score, evidence, and site-class warnings.",
                "Turn site calibration rows into before/after parser regression fixtures with confidence-delta reporting.",
                "Add R-group replay scenario comparison across source versions, with local rollback rehearsal packets.",
                "Use Reviewer Cockpit as the entry point for local notes, closure, and remediation; no procurement or real experiment feedback automation.",
            ],
        },
    ]


def build_pptx(slides_data: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for index, item in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PAPER if index else INK
        if index == 0:
            text_box(slide, 0.72, 0.62, 3.6, 0.3, "AUTOMEDCHEMIST", size=9, bold=True, color=RGBColor(124, 227, 211))
            text_box(slide, 0.72, 1.34, 8.2, 1.2, item["title"], size=38, bold=True, color=WHITE)
            text_box(slide, 0.76, 2.72, 8.4, 0.8, item["subtitle"], size=17, color=RGBColor(215, 225, 232))
            add_bullet_list(slide, 0.78, 4.35, 6.3, item["bullets"], size=15, color=WHITE)
            for i, label in enumerate(["native shell", "site selection", "governance", "exports"]):
                text_box(slide, 8.7, 1.1 + i * 0.78, 2.9, 0.36, label.upper(), size=15, bold=True, color=RGBColor(124, 227, 211))
            continue
        add_title(slide, item["title"])
        y = 1.75
        if item.get("metrics"):
            for col, metric in enumerate(item["metrics"]):
                add_metric(slide, 0.75 + col * 2.8, 1.65, metric[0], metric[1], [TEAL, GREEN, AMBER, CORAL][col % 4])
            y = 2.72
        add_bullet_list(slide, 0.85, y, 10.4, item["bullets"], size=17, color=INK)
        if item.get("metrics") and index in {3, 4, 5}:
            values = [metric[1] for metric in item["metrics"][1:]]
            numeric = [int(v) for v in values if isinstance(v, int)]
            max_value = max(numeric or [1])
            for offset, metric in enumerate(item["metrics"][1:]):
                if isinstance(metric[1], int):
                    add_bar(slide, 0.9, 5.55 + offset * 0.36, 5.2, metric[0], metric[1], max_value, [TEAL, AMBER, CORAL][offset % 3])
    prs.save(PPTX_OUT)


def font(size: int, bold: bool = False):
    candidates = ["segoeuib.ttf" if bold else "segoeui.ttf", "arialbd.ttf" if bold else "arial.ttf"]
    for name in candidates:
        try:
            return ImageFont.truetype(name, size)
        except Exception:
            continue
    return ImageFont.load_default()


def wrap(draw: ImageDraw.ImageDraw, text: str, max_width: int, fnt) -> list[str]:
    words = str(text).split()
    lines: list[str] = []
    current = ""
    for word in words:
        probe = f"{current} {word}".strip()
        if draw.textlength(probe, font=fnt) <= max_width:
            current = probe
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def build_previews(slides_data: list[dict]) -> list[Path]:
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    paths = []
    for index, item in enumerate(slides_data, start=1):
        image = Image.new("RGB", (1600, 900), "#17202A" if index == 1 else "#F6F8FA")
        draw = ImageDraw.Draw(image)
        title_font = font(58 if index == 1 else 42, bold=True)
        body_font = font(28)
        small_font = font(18, bold=True)
        x = 90
        y = 90 if index == 1 else 70
        draw.text((x, y), "AUTOMEDCHEMIST PRODUCT UPDATE", fill="#7CE3D3" if index == 1 else "#0F766E", font=small_font)
        y += 70
        for line in wrap(draw, item["title"], 1180, title_font):
            draw.text((x, y), line, fill="#FFFFFF" if index == 1 else "#17202A", font=title_font)
            y += 62
        if item.get("subtitle"):
            y += 18
            for line in wrap(draw, item["subtitle"], 1120, body_font):
                draw.text((x, y), line, fill="#D7E1E8", font=body_font)
                y += 38
        if item.get("metrics"):
            y += 26
            for col, metric in enumerate(item["metrics"]):
                mx = x + col * 330
                draw.text((mx, y), str(metric[0]).upper(), fill="#52616F", font=small_font)
                draw.text((mx, y + 30), str(metric[1]), fill=["#0F766E", "#1F7A52", "#B45309", "#BE4137"][col % 4], font=font(38, bold=True))
            y += 110
        y += 36
        fill = "#FFFFFF" if index == 1 else "#17202A"
        for bullet in item["bullets"]:
            lines = wrap(draw, bullet, 1260, body_font)
            draw.text((x, y), "-", fill=fill, font=body_font)
            for line in lines:
                draw.text((x + 34, y), line, fill=fill, font=body_font)
                y += 38
            y += 12
        path = PREVIEW_DIR / f"slide_{index:02d}.png"
        image.save(path)
        paths.append(path)
    sheet = Image.new("RGB", (1600, 900), "#F6F8FA")
    for idx, path in enumerate(paths):
        thumb = Image.open(path).resize((380, 214))
        x = 40 + (idx % 4) * 390
        y = 60 + (idx // 4) * 310
        sheet.paste(thumb, (x, y))
        ImageDraw.Draw(sheet).text((x, y + 222), f"Slide {idx + 1}", fill="#17202A", font=font(20, bold=True))
    sheet_path = PREVIEW_DIR / "contact_sheet.png"
    sheet.save(sheet_path)
    paths.append(sheet_path)
    return paths


def build_pdf(slides_data: list[dict]) -> None:
    pdf = canvas.Canvas(str(PDF_OUT), pagesize=landscape(letter))
    width, height = landscape(letter)
    for index, item in enumerate(slides_data, start=1):
        pdf.setFillColor(colors.HexColor("#17202A" if index == 1 else "#F6F8FA"))
        pdf.rect(0, 0, width, height, fill=True, stroke=False)
        pdf.setFillColor(colors.HexColor("#7CE3D3" if index == 1 else "#0F766E"))
        pdf.setFont("Helvetica-Bold", 8)
        pdf.drawString(0.55 * inch, height - 0.55 * inch, "AUTOMEDCHEMIST PRODUCT UPDATE")
        pdf.setFillColor(colors.white if index == 1 else colors.HexColor("#17202A"))
        pdf.setFont("Helvetica-Bold", 25)
        y = height - 1.15 * inch
        for line in item["title"].split("\n"):
            pdf.drawString(0.55 * inch, y, line[:80])
            y -= 0.36 * inch
        pdf.setFont("Helvetica", 12)
        y -= 0.25 * inch
        if item.get("metrics"):
            for col, metric in enumerate(item["metrics"]):
                x = 0.65 * inch + col * 1.95 * inch
                pdf.setFillColor(colors.HexColor("#52616F"))
                pdf.setFont("Helvetica-Bold", 7)
                pdf.drawString(x, y, str(metric[0]).upper())
                pdf.setFillColor(colors.HexColor(["#0F766E", "#1F7A52", "#B45309", "#BE4137"][col % 4]))
                pdf.setFont("Helvetica-Bold", 18)
                pdf.drawString(x, y - 0.28 * inch, str(metric[1]))
            y -= 0.9 * inch
        pdf.setFillColor(colors.white if index == 1 else colors.HexColor("#17202A"))
        pdf.setFont("Helvetica", 12)
        for bullet in item["bullets"]:
            pdf.drawString(0.7 * inch, y, f"- {bullet}"[:115])
            y -= 0.34 * inch
        pdf.showPage()
    pdf.save()


def main() -> None:
    slides_data = build_slide_data()
    build_pptx(slides_data)
    build_pdf(slides_data)
    previews = build_previews(slides_data)
    report = {
        "status": "built",
        "pptx": str(PPTX_OUT),
        "pdf": str(PDF_OUT),
        "preview_count": len(previews),
        "preview_paths": [str(path) for path in previews],
    }
    (ROOT / "data/releases/product_update_deck_report.json").write_text(
        json.dumps(report, indent=2, sort_keys=True),
        encoding="utf-8",
    )
    print(json.dumps(report, indent=2, sort_keys=True))


if __name__ == "__main__":
    main()
